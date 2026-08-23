"""The pptreport example, built and read back.

``powerpoint_input.yaml`` drives every kind of fill the pptx back end has:
layout copies (``TitleSlide``, ``TaskProjectDefinition``, ... ``BackCover``),
placeholder texts, ``_global_`` fills that reach every slide (reference, date
and name), parameter-file values, pictures with and without a description, a
table from a CSV file, two videos with poster frames, and two pictures whose
files do not exist. A test that only checks the file is there proves none of
that, so this module reads the deck back:

* what it *says*, against the reference in ``expected/powerpoint_input.json``
  -- captured at `44267a8` from the ``.rdf`` this fixture was translated from,
  by the differential harness, which this case has graduated from (see
  ``02_basetest/differential``; the reading and normalising are shared in
  ``common_case``);
* what it *shows*, which no text comparison can see: the pictures with their
  sizes, the table with its shape, the two movies with their posters, and the
  two pictures that could only be announced as missing.
"""

from pathlib import Path
import importlib.util
import os
import sys

import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_MEDIA_TYPE
from pptx.util import Cm

THIS_DIR = Path(__file__).resolve().parent
CASE_ROOT = Path(__file__).resolve().parent.parent
if str(CASE_ROOT) not in sys.path:
    sys.path.append(str(CASE_ROOT))

from _setup_examples import *

module_path = Path(__file__).resolve().parent.parent.parent / '02_basetest' / 'common_case.py'

# Load the module from the given path
spec = importlib.util.spec_from_file_location('common_case', str(module_path))
module = importlib.util.module_from_spec(spec)
spec.loader.exec_module(module)

CaseConfig = module.CaseConfig
run_pptx_case = module.run_pptx_case
said, normalise, reference, difference = (
    module.said, module.normalise, module.reference, module.difference)

REFERENCE = THIS_DIR / 'expected' / 'powerpoint_input.json'


def build(tmp_path):
    """The deck, typeset the way every pptx case is."""
    config = CaseConfig(
        name="pptreport",
        case_dir=THIS_DIR,
        document_name="powerpoint_input.yaml",
        template_doc_name="template.pptx",
        output_name="final_report.pptx",
        include_patterns=["*.yaml", "template.pptx"],
        data_source_dir=DATA_SOURCE,
        finish=False,
        createpdf=False,
    )
    return run_pptx_case(config, tmp_path)


def portable(lines, workspace):
    """*lines* with the workspace taken out of the paths quoted in them.

    The runner hands ``ReportDataFile`` an absolute document path, so the data
    directory -- and the path a missing-file message quotes -- is absolute,
    while the reference was captured from a run with a relative one and reads
    ``data/...``. The message quotes with ``repr``, which on Windows doubles
    every backslash, so the prefix to strip is the repr'd workspace; the
    separators are folded to ``/`` afterwards (on both sides, see the
    comparison) so the test reads the same on every platform.
    """
    prefix = repr(str(workspace) + os.sep)[1:-1]
    return [fold(line.replace(prefix, '')) for line in lines]


def fold(line):
    """A repr'd Windows path reads like a POSIX one."""
    return line.replace('\\\\', '/')


def of_kind(slide, kind):
    return [shape for shape in slide.shapes if shape.shape_type == kind]


def size(shape):
    """(width, height) in cm, as the slide shows it."""
    return (round(shape.width / Cm(1), 2), round(shape.height / Cm(1), 2))


def test_document_is_created(tmp_path):
    print(f'\nWorking in {tmp_path}')

    result_path = build(tmp_path)

    assert result_path.exists(), "Expected final_report.pptx to be generated"
    assert result_path.stat().st_size > 0, "Generated document should not be empty"


def test_the_deck_says_what_the_reference_says(tmp_path):
    """Slide by slide, every text the deck carries, against the stored
    reference -- digits and weekday names collapsed on both sides, since the
    reference was captured on another day and ``date: now`` is evaluated per
    run. The one line that names a path (a picture whose file is missing) is
    made portable first, see :func:`portable`."""
    deck = build(tmp_path)

    got = normalise(portable(said(deck), deck.parent))
    expected = [fold(line) for line in reference(REFERENCE)]

    assert got == expected, difference(expected, got)


def test_pictures_tables_and_videos_are_placed(tmp_path):
    """What the text comparison cannot see. The layout sequence is the
    ``_content_`` sequence (the template's own first slide removed); the
    sizes are what python-pptx computes from each file and the tag, so a
    picture that silently went missing, or one placed at native size instead
    of the tag's, shows up here."""
    deck = pptx.Presentation(build(tmp_path))
    slides = list(deck.slides)

    assert [slide.slide_layout.name for slide in slides] == [
        'TitleSlide', 'TaskProjectDefinition', 'Material',
        'ResultsGeneral', 'ResultsGeneral', 'ResultsGeneral',
        'ResultsAnimation', 'ResultsTwoBlocks', 'BackCover']
    title, definition, material, empty, result, video, animation, two_blocks, back = slides

    # pictures: image:main_model and image:model_icon (screw.png, square),
    # image:history (curve_test.png) -- one each, at the size the slide shows
    pictures = lambda slide: [size(p) for p in of_kind(slide, MSO_SHAPE_TYPE.PICTURE)]
    assert pictures(title) == [(7.99, 7.99)]
    assert pictures(definition) == [(3.4, 3.4)]
    assert pictures(result) == [(12.58, 8.39)]
    assert pictures(empty) == [], 'an empty ResultsGeneral entry adds a bare slide'
    assert pictures(two_blocks) == [], 'both of its files are missing'

    # the table from table3.csv, with its description as a caption
    (table,) = of_kind(material, MSO_SHAPE_TYPE.TABLE)
    assert (len(table.table.rows), len(table.table.columns)) == (4, 5)
    assert table.table.cell(0, 0).text == 'Type'
    captions = [shape.text_frame.text for shape in material.shapes
                if shape.has_text_frame and shape.text_frame.text.startswith('Table:')]
    assert captions == ['Table: Speed records of cars']

    # the two videos, each with its poster frame
    for slide in (video, animation):
        (movie,) = of_kind(slide, MSO_SHAPE_TYPE.MEDIA)
        assert movie.media_type == PP_MEDIA_TYPE.MOVIE
        assert movie.poster_frame.content_type == 'image/jpeg'

    # the two missing pictures are announced, not shown
    announced = [shape.text_frame.text for shape in two_blocks.shapes
                 if shape.has_text_frame and 'non existing image file' in shape.text_frame.text]
    assert len(announced) == 2
    assert all('result1' in text for text in announced)

    assert len(back.shapes) == 0, 'BackCover takes nothing from the document'
    assert deck.core_properties.author.startswith('Scriptum ')
