"""Reusable helpers for docx and pptx creation tests.

This module extracts the reusable pieces from the legacy notebook-based
"CreateDOCforEssay" test so that new scenarios can reuse the same setup and
execution flow. To create a new case, define a :class:`CaseConfig` pointing at
the report document (``.yaml``) and the template and call
:func:`run_docx_case` or :func:`run_pptx_case`.

Reading the result back
-----------------------
A case test that only asserts a file exists proves nothing about the document
(the clone-ordering defect fixed in `b1a4afd` lived for years behind such
tests). The second half of this module reads a finished document back so a
case can compare what it *says* with the reference kept beside it in
``expected/<stem>.json``: :func:`said` lists the texts of a ``.docx`` or a
``.pptx`` in order, :func:`normalise` hides what changes from run to run
(digits, weekday names -- both for dates), :func:`reference` reads a stored
list through the same normaliser, and :func:`difference` says where two lists
part. Two more make a run comparable with a reference captured another way:
:func:`portable` takes the absolute workspace out of a quoted path (the
runner hands ``ReportDataFile`` an absolute document path, a reference was
captured with a relative one), and :func:`comparable` drops the field results
only Word refreshes. Every case test uses these, so no two tests can disagree
about what "the same document" means.

What a document *shows* is read back too, since no text comparison sees a
picture that silently went missing (`1f75367`): :func:`drawings` lists the
inline pictures of a Word paragraph with their sizes, :func:`shapes` and
:func:`size_cm` do the same for the shapes of a slide.
"""

import json
import os
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable, Sequence

import docx
from docx.oxml.ns import qn
import pptx

from _setup_basetest import *
import Scriptum  # type: ignore



@dataclass
class CaseConfig:
    """Configuration for a document generation scenario.

    Attributes:
        name: Identifier used for debugging output.
        case_dir: Directory that contains the report document and the template.
        document_name: Name of the ``.yaml`` report document in ``case_dir``.
        template_doc_name: Name of the input Word/PowerPoint document.
        output_name: Name of the generated document.
        include_patterns: Glob patterns (relative to ``case_dir``) that should
            be linked into the workspace before running the test.
        data_source_dir: Optional override for the shared data source folder.
        finish: open in Word/PowerPoint to update tables etc
        createpdf: save from Word/PowerPoint as PDF
    """

    name: str
    case_dir: Path
    document_name: str
    template_doc_name: str
    output_name: str
    include_patterns: Sequence[str] = field(default_factory=list)
    data_source_dir: Path | None = None
    finish: bool = False
    createpdf: bool = False


class WorkspaceBuilder:
    """Prepare a disposable workspace for document generation tests."""

    def __init__(self, tmp_path: Path, data_source_dir: Path = DATA_SOURCE):
        self.tmp_path = tmp_path
        self.data_source_dir = data_source_dir

    def build(self, case_dir: Path, include_patterns: Iterable[str]) -> Path:
        """Create a workspace with linked input data.

        The workspace mirrors the steps from the original notebook: it links the
        shared ``data_source`` folder and any case-specific files that match the
        provided glob patterns.
        """

        workdir = self.tmp_path / "workspace"
        workdir.mkdir()

        ensure_link(self.data_source_dir, workdir / "data")

        include_files: set[Path] = set()
        for pattern in include_patterns:
            include_files.update(case_dir.glob(pattern))

        if not include_files:
            patterns = ", ".join(include_patterns)
            msg = f"No files found in {case_dir} for patterns: {patterns}"
            raise FileNotFoundError(msg)

        for src in include_files:
            ensure_link(src, workdir / src.name)

        print(f'WORKSPACE: {workdir}')

        return workdir


def run_docx_case(config: CaseConfig, tmp_path: Path) -> Path:
    """Execute a document generation test based on the provided configuration.

    The function handles workspace creation, running Scriptum to typeset the
    document, and returning the path to the generated file.
    """

    workspace = WorkspaceBuilder(
        tmp_path, config.data_source_dir or DATA_SOURCE
    ).build(config.case_dir, config.include_patterns)

    current_dir = Path(os.getcwd())
    os.chdir(workspace)
    try:
        rdf = Scriptum.ReportDataFile(workspace / config.document_name)

        document = Scriptum.ManagedDocx(config.template_doc_name)
        document.typesetting(
            rdf,
            addcopy=True,
            directfill=True,
            globalfill=True,
            cleanup=True,
            removetemplate=True,
            cleardust=True,
            setproperties=True,
        )

        output_path = workspace / config.output_name
        document.save(config.output_name, finish=config.finish, createpdf=config.createpdf)
    finally:
        os.chdir(current_dir)

    return output_path

def run_pptx_case(config: CaseConfig, tmp_path: Path) -> Path:
    """Execute a document generation test based on the provided configuration.

    The function handles workspace creation, running Scriptum to typeset the
    document, and returning the path to the generated file.
    """

    workspace = WorkspaceBuilder(
        tmp_path, config.data_source_dir or DATA_SOURCE
    ).build(config.case_dir, config.include_patterns)

    current_dir = Path(os.getcwd())
    os.chdir(workspace)
    try:
        rdf = Scriptum.ReportDataFile(workspace / config.document_name)

        document = Scriptum.ManagedPptx(config.template_doc_name)
        document.artist(
            rdf,
            directfill=True,
            globalfill=True,
            cleardust=True,
            setproperties=True,
        )

        output_path = workspace / config.output_name
        document.document.core_properties.title='AutoReport'
        document.remove_slide(0)
    
        document.save(config.output_name, finish=config.finish, createpdf=config.createpdf)
    finally:
        os.chdir(current_dir)

    return output_path


# ------------------------------------------------------------ reading back

DIGITS = re.compile(r'\d+')
#: The default datetime format used to start with a weekday name, which no
#: amount of digit-collapsing hides: a reference is captured on another day.
WEEKDAY = re.compile(r'\b(?:Mon|Tue|Wed|Thu|Fri|Sat|Sun)\b')


def said(path) -> list:
    """What a finished document says, in order: the non-empty texts.

    A ``.docx``: every paragraph, then every table cell. A ``.pptx``: slide by
    slide, every paragraph of every shape with a text frame, then the cells of
    every table shape. Whitespace is stripped and empty texts are dropped, so
    a placeholder a run left blank does not count as something said.
    """
    path = Path(path)
    lines = _slides(path) if path.suffix.lower() == '.pptx' else _paragraphs(path)
    return [line for line in lines if line]


def _paragraphs(path):
    document = docx.Document(path)
    lines = [p.text.strip() for p in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            lines.extend(cell.text.strip() for cell in row.cells)
    return lines


def _slides(path):
    lines = []
    for slide in pptx.Presentation(path).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    lines.append(''.join(run.text for run in paragraph.runs).strip())
            if shape.has_table:
                for row in shape.table.rows:
                    lines.extend(cell.text.strip() for cell in row.cells)
    return lines


def portable(lines, workspace) -> list:
    """*lines* with *workspace* taken out of the paths quoted in them.

    The runner hands ``ReportDataFile`` an absolute document path, so the data
    directory -- and the path a message quotes for a missing picture or video
    -- is absolute, while a reference was captured from a run with a relative
    one and reads ``data/...``. The message quotes with ``repr``, which on
    Windows doubles every backslash, so the prefix to strip is the repr'd
    workspace; the separators are then folded to ``/`` so the comparison reads
    the same on every platform. Apply :func:`fold` to the reference side too.
    """
    prefix = repr(str(workspace) + os.sep)[1:-1]
    return [fold(line.replace(prefix, '')) for line in lines]


def fold(line) -> str:
    """A repr'd Windows path reads like a POSIX one."""
    return line.replace('\\\\', '/')


#: A field result of a list of tables / figures: ``caption<TAB>page``, after
#: the digits have been collapsed. Only Word updates these, so a finished run
#: and a plain run disagree on them by construction; see :func:`comparable`.
FIELD_ENTRY = re.compile(r'\t#$')


def comparable(lines) -> list:
    """*lines* without the field results only Word updates.

    A Word template carries a *list of tables* and a *list of figures* -- TOC
    fields whose stored result is whatever Word last wrote into them. A plain
    run (``save()`` without ``finish``) leaves that result as the template had
    it; a ``finish=True`` run -- Windows with Word, the only thing that can
    update a field -- rewrites it with one entry per caption actually in the
    document. So the same document says two different things in those lines
    depending on where it was built, and a reference captured one way can
    never match a run made the other way. The entries are the one shape
    ``caption<TAB>page`` and are dropped from **both** sides; nothing is lost,
    every caption is also a caption paragraph of its own, which stays compared.
    The wordreport reference was captured with ``finish=True`` (``e60f4e0``),
    the others without, and all compare against a plain run anywhere and a
    finished run where Word is.
    """
    return [line for line in lines if not FIELD_ENTRY.search(line)]


def normalise(lines) -> list:
    """*lines* with runs of digits and weekday names collapsed to ``#``.

    Both are for dates: a document using ``date: now`` is evaluated when it is
    built, and the reference was built on another day. What a comparison
    exists for -- a text missing, an extra one, two in the wrong order --
    survives it.
    """
    return [WEEKDAY.sub('#', DIGITS.sub('#', line)) for line in lines]


def reference(path) -> list:
    """The stored reference at *path* (a JSON list of texts), normalised.

    It was captured on another day, so it carries that day's digits and
    weekday name -- normalising only one side of a comparison is how you end
    up measuring the calendar.
    """
    return normalise(json.loads(Path(path).read_text(encoding='utf-8')))


def difference(expected, got) -> str:
    """A readable account of the first place the two diverge."""
    for index, (a, b) in enumerate(zip(expected, got)):
        if a != b:
            return (f'first difference at line {index}:\n'
                    f'  expected: {a[:120]!r}\n'
                    f'  got     : {b[:120]!r}')
    if len(expected) != len(got):
        longer, name = ((expected, 'the reference') if len(expected) > len(got)
                        else (got, 'this run'))
        extra = longer[min(len(expected), len(got)):][:4]
        return (f'{name} says {abs(len(expected) - len(got))} more line(s): '
                f'{[line[:60] for line in extra]}')
    return 'no difference'


# ----------------------------------------------------- what a document shows

#: Both python-docx and python-pptx measure in EMU; sizes are read back in cm.
EMU_PER_CM = 360000


def drawings(paragraph) -> list:
    """(width, height) in cm of every inline picture in a Word *paragraph*."""
    return [(round(int(extent.get('cx')) / EMU_PER_CM, 2),
             round(int(extent.get('cy')) / EMU_PER_CM, 2))
            for extent in paragraph._p.iter(qn('wp:extent'))]


def shapes(slide, kind) -> list:
    """The shapes of a *slide* of one ``MSO_SHAPE_TYPE`` *kind*, in order."""
    return [shape for shape in slide.shapes if shape.shape_type == kind]


def size_cm(shape) -> tuple:
    """(width, height) in cm of a slide shape, as the slide shows it."""
    return (round(shape.width / EMU_PER_CM, 2), round(shape.height / EMU_PER_CM, 2))
