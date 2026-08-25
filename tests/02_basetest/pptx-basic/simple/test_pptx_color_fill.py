"""A ``color`` modifier paints the font of the text its fill writes.

The first thing a colour actually *does* in a document: the loader has long
parsed and diagnosed colours (``values/test_color_value.py``, the yaml_loader
tests) while both back ends dropped them -- a direct ``color:`` fill erased
its tag and a ``color`` modifier was ignored. Now the run that receives the
fill's text is painted, and nothing else is. The docx twin is
``docx_basic/simple/test_docx_color_fill.py``; the template-*add* route (a
``p:default`` clone carrying ``color: green``) is pinned by the pptreport
example, whose fixture has said ``color: green`` all along without anything
reading it.
"""

from pathlib import Path
import shutil
import sys

import pptx
from pptx.dml.color import RGBColor

THIS_DIR = Path(__file__).resolve().parent
CASE_ROOT = THIS_DIR.parent
if str(CASE_ROOT) not in sys.path:
    sys.path.append(str(CASE_ROOT))

from _setup_pptx_basic import *          # noqa: F401,F403

DOCUMENT = """_scriptum_:
  version: 4
  documenttype: pptx
  datadir: ./data
_content_:
  - TitleSlide:
      - title:
          text: Corporate Report
          color: '005596'
      - subtitle:
          text: DRAFT
          color: crimson
"""


def test_a_color_modifier_paints_the_placeholder_text(tmp_path, monkeypatch):
    import Scriptum

    shutil.copy(THIS_DIR / 'template.pptx', tmp_path)
    (tmp_path / 'data').mkdir()
    (tmp_path / 'case.yaml').write_text(DOCUMENT, encoding='utf-8')
    monkeypatch.chdir(tmp_path)

    rdf = Scriptum.ReportDataFile('case.yaml')
    deck = Scriptum.ManagedPptx('template.pptx')
    deck.artist(rdf, directfill=True, globalfill=True,
                cleardust=True, setproperties=True)
    deck.remove_slide(0)
    deck.save('out.pptx', finish=False, createpdf=False)

    colors = {run.text: run.font.color
              for slide in pptx.Presentation('out.pptx').slides
              for shape in slide.shapes if shape.has_text_frame
              for paragraph in shape.text_frame.paragraphs
              for run in paragraph.runs}

    assert colors['Corporate Report'].rgb == RGBColor(0x00, 0x55, 0x96)
    assert colors['DRAFT'].rgb == RGBColor(0xDC, 0x14, 0x3C)   # crimson
