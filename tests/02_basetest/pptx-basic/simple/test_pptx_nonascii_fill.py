"""Umlauts, accents and emoji pass through a slide fill unchanged.

The pptx twin of ``docx_basic/simple/test_docx_nonascii_fill.py``: a value
travels document -> loader -> task -> placeholder run, and nothing on the way
may re-encode it. The document title takes the same trip into the deck's core
properties.
"""

from pathlib import Path
import shutil
import sys

import pptx

THIS_DIR = Path(__file__).resolve().parent
CASE_ROOT = THIS_DIR.parent
if str(CASE_ROOT) not in sys.path:
    sys.path.append(str(CASE_ROOT))

from _setup_pptx_basic import *          # noqa: F401,F403

DOCUMENT = """_scriptum_:
  version: 4
  documenttype: pptx
  datadir: ./data
  documenttitle: 'Prüfbericht 🚀'
_content_:
  - TitleSlide:
      - title: 'Größenprüfung äöü ß'
      - subtitle: 'Résumé ✓ 100 € — naïve façade'
"""


def test_nonascii_values_reach_the_deck_verbatim(tmp_path, monkeypatch):
    import Scriptum

    shutil.copy(THIS_DIR / 'template.pptx', tmp_path)
    (tmp_path / 'data').mkdir()
    (tmp_path / 'case.yaml').write_text(DOCUMENT, encoding='utf-8')
    monkeypatch.chdir(tmp_path)

    rdf = Scriptum.ReportDataFile('case.yaml')
    deck = Scriptum.ManagedPptx('template.pptx')
    deck.artist(rdf, directfill=True, globalfill=True,
                cleardust=True, setproperties=True)
    deck.remove_slide(0)
    deck.save('out.pptx', finish=False, createpdf=False)

    finished = pptx.Presentation('out.pptx')
    said = [''.join(run.text for run in paragraph.runs)
            for slide in finished.slides
            for shape in slide.shapes if shape.has_text_frame
            for paragraph in shape.text_frame.paragraphs]

    assert any('Größenprüfung äöü ß' in text for text in said), said
    assert any('Résumé ✓ 100 € — naïve façade' in text for text in said), said
    assert finished.core_properties.title == 'Prüfbericht 🚀'
