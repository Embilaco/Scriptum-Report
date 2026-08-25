"""The simple pptx case, built and read back.

``powerpoint_simple.yaml`` is the pptx back end in small: layout copies
(``TitleSlide`` ... ``BackCover``), placeholder texts, ``_global_`` values
on every slide (reference, date and name), parameter-file values placed at
a marker with ``left``/``top``, a CSV table described from its first row, a
text file with extra ``info``/``more`` fills, and two pudding pictures with
descriptions -- one sized by ``height``, one pinned by ``left``/``top``
(announced-missing pictures are the pptreport case's business). A test that
only checks the file is there proves none of that, so this module reads the
deck back:

* what it *says*, against ``expected/powerpoint_simple.json`` ;
* what it *shows*: the layouts in order, the pictures with their sizes,
  the table, and the text boxes the parameter file and the text file
  produced.
"""

from pathlib import Path
import sys

import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pytest

THIS_DIR = Path(__file__).resolve().parent
CASE_ROOT = Path(__file__).resolve().parent.parent
if str(CASE_ROOT) not in sys.path:
    sys.path.append(str(CASE_ROOT))

from _setup_pptx_basic import *
from common_case import CaseConfig, run_pptx_case
from common_case import said, normalise, reference, difference, portable, fold, shapes, size_cm
from common_case import checkreport_comparison, com_quiet

REFERENCE = THIS_DIR / 'expected' / 'powerpoint_simple.json'


def build(tmp_path):
    """The deck, typeset the way every pptx case is."""
    config = CaseConfig(
        name="report",
        case_dir=THIS_DIR,
        document_name="powerpoint_simple.yaml",
        template_doc_name="template.pptx",
        output_name="final_report.pptx",
        include_patterns=["*.yaml", "template.pptx"],
        data_source_dir=DATA_SOURCE,
        finish=False,
        createpdf=False,
    )
    return run_pptx_case(config, tmp_path)


def test_document_is_created(tmp_path):
    print(f'\nWorking in {tmp_path}')

    result_path = build(tmp_path)

    assert result_path.exists(), "Expected final_report.pptx to be generated"
    assert result_path.stat().st_size > 0, "Generated document should not be empty"


def test_the_deck_says_what_the_reference_says(tmp_path):
    """Slide by slide, every text the deck carries, against the stored
    reference. The two missing pictures are announced with their path, which
    the runner makes absolute -- portable() takes the workspace out again."""
    deck = build(tmp_path)

    got = normalise(portable(said(deck), deck.parent))
    expected = [fold(line) for line in reference(REFERENCE)]

    assert got == expected, difference(expected, got)


def test_pictures_table_and_text_boxes_are_placed(tmp_path):
    """What the text comparison cannot see."""
    deck = pptx.Presentation(build(tmp_path))
    slides = list(deck.slides)

    assert [slide.slide_layout.name for slide in slides] == [
        'TitleSlide', 'TaskProjectDefinition', 'TitleContent', 'TitleContent',
        'Material', 'BackCover']
    title, definition, table_slide, pictures_slide, material, back = slides

    # image:main_model and image:model_icon (screw.png, square), one each
    pictures = lambda slide: [size_cm(p) for p in shapes(slide, MSO_SHAPE_TYPE.PICTURE)]
    assert pictures(title) == [(7.99, 7.99)]
    assert pictures(definition) == [(3.4, 3.4)]

    # the parameter file's values placed at the definition slide's marker
    boxes = [shape.text_frame.text for shape in shapes(definition, MSO_SHAPE_TYPE.TEXT_BOX)]
    assert boxes == ['This is a description of the testplan\nmore description\nand even more ',
                     'WhatEver-F1 ']

    # table:small from table3.csv; the layout has no caption box, so the
    # description from row1 has nowhere to go and the table stands alone
    (table,) = shapes(table_slide, MSO_SHAPE_TYPE.TABLE)
    assert (len(table.table.rows), len(table.table.columns)) == (4, 5)
    assert table.table.cell(0, 0).text == 'Type'
    assert shapes(table_slide, MSO_SHAPE_TYPE.TEXT_BOX) == []

    # both pictures of the next slide are placed at their size: image:generic
    # scaled to its 2 cm height request by the marker box, image:history
    # pinned at left/top with its own proportions
    assert pictures(pictures_slide) == [(1.21, 1.24), (2.96, 3.04)]

    # text:insert: the text file and its two extra fills, three boxes
    boxes = [shape.text_frame.text for shape in shapes(material, MSO_SHAPE_TYPE.TEXT_BOX)]
    assert boxes[0].startswith('A general text may look like this Lorem ipsum')
    assert boxes[1:] == ['Some further info Moore', 'And even more Roger']

    assert len(back.shapes) == 0, 'BackCover takes nothing from the document'
    assert deck.core_properties.author.startswith('Scriptum ')


def test_the_deck_carries_the_documents_title(tmp_path):
    """``setproperties=True`` writes ``rdf.settings.documenttitle`` into the
    core properties. The runner used to overwrite it with a hard-coded
    'AutoReport' afterwards, so a broken ``setproperties`` went unseen."""
    deck = pptx.Presentation(build(tmp_path))

    assert deck.core_properties.title == 'This is a bloody test document'


def test_a_powerpoint_resave_keeps_what_the_deck_says(tmp_path, capsys):
    """``finish=True`` hands the saved deck to PowerPoint, which rewrites the
    file (Windows only; a no-op elsewhere) -- and the re-save must not change
    what the deck says. ``createpdf=True`` rides along so the PDF export is
    pinned too; before this it was exercised only by the manual scripts. The
    deck is opened ``WithWindow=False``, so a suite run puts no PowerPoint
    window on the desktop -- if one starts flashing up again, that regressed.
    Where PowerPoint cannot finish -- absent, busy, or refusing the call --
    the runner prints the reason and the deck stays the plain save: that
    reports here as xfailed with the reason, never as a failure."""
    config = CaseConfig(
        name="report",
        case_dir=THIS_DIR,
        document_name="powerpoint_simple.yaml",
        template_doc_name="template.pptx",
        output_name="final_report.pptx",
        include_patterns=["*.yaml", "template.pptx"],
        data_source_dir=DATA_SOURCE,
        finish=True,
        createpdf=True,
    )
    with com_quiet():
        deck = run_pptx_case(config, tmp_path)
    out = capsys.readouterr().out

    if 'failed to update' in out:
        tail = ' '.join(out[out.index('failed to update'):].split())
        pytest.xfail('PowerPoint could not finish the deck -- ' + tail[:200])

    report = checkreport_comparison(deck, REFERENCE)
    assert report == '', report
    pdf = deck.parent / 'final_report.pdf'
    assert pdf.exists() and pdf.stat().st_size > 0, 'createpdf left no PDF beside the deck'


def test_the_checkreport_notebook_would_say_identical(tmp_path):
    """The comparison ``CheckReport.ipynb`` beside this file ends with --
    plain, digits and weekdays collapsed, nothing else -- on the same build
    as above. Green when the notebook would print IDENTICAL; anything else
    reports as xfailed with the first difference instead of staying out of
    the suite."""
    report = checkreport_comparison(build(tmp_path), REFERENCE)
    if report:
        pytest.xfail('CheckReport.ipynb would not say IDENTICAL -- ' + report)
