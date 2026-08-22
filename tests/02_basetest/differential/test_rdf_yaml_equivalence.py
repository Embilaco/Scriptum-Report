"""What each fixture's ``.yaml`` generates, against what its ``.rdf`` did.

This is the safety net for the last piece of the migration. Every docx case
test asserts only that a file was written and is not empty, which is how the
clone-ordering defect fixed in `b1a4afd` lived in shipped code for years with a
green suite. This reads the finished documents back instead.

Why a stored reference rather than a live comparison
----------------------------------------------------
It began by generating both sides in one run. That stopped working the moment
``StructuredElement.path`` became canonical: the text parser emits
``section:title`` where the tree now says ``section:title::1``, so the *old*
side degrades and there is nothing left to compare against.

So the reference is captured instead. Each case keeps its own beside the
fixture it belongs to -- ``<case>/expected/<stem>.json`` -- which is what each
``.rdf`` generated at `44267a8`, the last commit before the tree changed. That
decouples the check from a parser which is on its way out: when the text format
goes, these keep working.

Living next to the fixture rather than in one pile means a case is a directory
you can read end to end: the document, the template, the data and what it is
supposed to say.

What is compared
----------------
The sequence of non-empty paragraph texts, then every table cell, with runs of
digits **and weekday names** collapsed to ``#``. Both are for dates: a fixture
using ``date:now`` is evaluated when the document is built, and the reference
was built on another day -- which the digits hide and the leading ``Fri`` of
the default format does not. What the comparison exists for -- a paragraph
missing, an extra one, two in the wrong order -- survives it.

One trap this had to avoid
--------------------------
The reader keeps **process-global state** and one root document per interpreter
is the standing rule, so :func:`reset` clears it before each generation rather
than relying on the suite's autouse fixture. A previous session was misled by
exactly that.
"""

from __future__ import annotations

import contextlib
import io
import json
import os
import re
import shutil
import tempfile
from pathlib import Path

import docx
import pptx
import pytest

import Scriptum
from Scriptum.rdf.reportDataFile import ReportDataFile
from Scriptum.rdf.tasks import ReportTask

TESTS_ROOT = Path(__file__).resolve().parents[2]
DIGITS = re.compile(r'\d+')
#: The default datetime format starts with a weekday name, which no amount
#: of digit-collapsing hides: the reference was captured on another day.
WEEKDAY = re.compile(r'\b(?:Mon|Tue|Wed|Thu|Fri|Sat|Sun)\b')

#: (case directory relative to tests/, fixture stem, template file)
CASES = [
    ('02_basetest/docx_basic/simple', 'word_simple', 'template.docx'),
    ('02_basetest/docx_basic/images', 'word_images', 'template_image.docx'),
    ('02_basetest/docx_basic/tables', 'word_tables', 'template_table.docx'),
    ('04_examples/wordreport', 'word_input', 'template.docx'),
    ('04_examples/essay', 'essay', 'essay.docx'),
    ('02_basetest/pptx-basic/simple', 'powerpoint_simple', 'template.pptx'),
    ('04_examples/pptreport', 'powerpoint_input', 'template.pptx'),
]

#: The one case that does not match, and why. It is a template problem, not a
#: loader one: ``template_text.docx`` spells its depth-3 block
#: ``<subsubsubsection:secondsubsubi>``, which is not in the docx ladder --
#: every other template in the corpus stops at ``subsubsection``. The ``.yaml``
#: uses the ladder's name, ``sub3section``, so the template needs the matching
#: rename before the two can agree.
PENDING = [
    ('02_basetest/docx_basic/text', 'word_text', 'template_text.docx'),
]

TEMPLATE_MISMATCH = (
    'template_text.docx spells a depth-3 block <subsubsubsection:...>, which '
    'is not in the docx ladder; the .yaml uses the ladder name sub3section, so '
    'the template needs the matching tag rename'
)


def reset():
    """Clear the process-global state the reader keeps between generations."""
    ReportTask._serial = 0
    ReportTask._tree = {}
    ReportTask._allPaths = {}
    ReportTask._newPaths = {}
    ReportDataFile._depth = 0
    ReportDataFile._global_settings = {}


def prepare(case):
    """A directory holding the fixtures, the template and a ``data``."""
    work = Path(tempfile.mkdtemp())
    source = TESTS_ROOT / case

    # Files only, so `expected/` stays where it is: the reference is what the
    # run is checked against, not an input to it.
    for pattern in ('*.yaml', '*.docx', '*.pptx'):
        for path in source.glob(pattern):
            shutil.copy(path, work)

    own_data = source / 'data'
    shutil.copytree(own_data if own_data.is_dir() else TESTS_ROOT / 'data_source',
                    work / 'data', dirs_exist_ok=True)
    return work


def generate(work, document_name, template, output=None):
    """Build one document. Returns what the run printed, for diagnosis.

    The two back ends are driven differently and always have been: Word
    typesets a document that already exists, PowerPoint *assembles* one from
    layouts and then drops the placeholder slide it started from. The steps
    here mirror ``common_case.run_docx_case`` and ``run_pptx_case``.
    """
    powerpoint = template.endswith('.pptx')
    output = output or ('out.pptx' if powerpoint else 'out.docx')

    os.chdir(work)
    reset()
    with contextlib.redirect_stdout(io.StringIO()) as printed:
        rdf = ReportDataFile(document_name)
        if powerpoint:
            managed = Scriptum.ManagedPptx(template)
            managed.artist(rdf, directfill=True, globalfill=True,
                           cleardust=True, setproperties=True)
            managed.remove_slide(0)
        else:
            managed = Scriptum.ManagedDocx(template, rdf)
            managed.typesetting(rdf)
        managed.save(output)
    return printed.getvalue()


def spoken(path):
    """What the finished document says, in order, dates neutralised."""
    said = _slides(path) if str(path).endswith('.pptx') else _paragraphs(path)
    return [WEEKDAY.sub('#', DIGITS.sub('#', line)) for line in said if line]


def _paragraphs(path):
    document = docx.Document(path)
    said = [p.text.strip() for p in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            said.extend(cell.text.strip() for cell in row.cells)
    return said


def _slides(path):
    said = []
    for slide in pptx.Presentation(path).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    said.append(''.join(run.text for run in paragraph.runs).strip())
            if shape.has_table:
                for row in shape.table.rows:
                    said.extend(cell.text.strip() for cell in row.cells)
    return said


def reference_path(case, stem):
    """Where a case keeps what it is supposed to say."""
    return TESTS_ROOT / case / 'expected' / f'{stem}.json'


def reference(case, stem):
    """The stored reference, through the same normaliser as a fresh run.

    It was captured on another day, so it carries that day's weekday name --
    and normalising only one side of a comparison is how you end up measuring
    the calendar.
    """
    stored = json.loads(reference_path(case, stem).read_text(encoding='utf-8'))
    return [WEEKDAY.sub('#', DIGITS.sub('#', line)) for line in stored]


def difference(expected, got):
    """A readable account of the first place the two diverge."""
    for index, (a, b) in enumerate(zip(expected, got)):
        if a != b:
            return (f'first difference at line {index}:\n'
                    f'  expected: {a[:120]!r}\n'
                    f'  got     : {b[:120]!r}')
    if len(expected) != len(got):
        longer, name = ((expected, 'the reference') if len(expected) > len(got)
                        else (got, 'this run'))
        extra = longer[min(len(expected), len(got)):][:4]
        return (f'{name} says {abs(len(expected) - len(got))} more line(s): '
                f'{[line[:60] for line in extra]}')
    return 'no difference'


def compare(case, stem, template):
    work = prepare(case)
    printed = generate(work, f'{stem}.yaml', template)
    got = spoken(work / ('out.pptx' if template.endswith('.pptx') else 'out.docx'))
    expected = reference(case, stem)
    complaints = [line for line in printed.splitlines()
                  if 'ERROR' in line or 'WARNING' in line]
    return expected, got, complaints


# ------------------------------------------------------- the harness works

def test_every_case_has_a_reference():
    """A missing reference would make a comparison vacuous."""
    for case, stem, template in CASES + PENDING:
        assert reference_path(case, stem).is_file(), \
            f'{stem}: no reference at {reference_path(case, stem)}'
        assert len(reference(case, stem)) > 3, stem


def test_no_reference_is_left_behind():
    """Every ``expected/`` in the tree belongs to a case listed here.

    A reference nobody compares against is worse than none: it looks like
    coverage and is not.
    """
    on_disk = {path.resolve() for path in TESTS_ROOT.rglob('expected/*.json')}
    claimed = {reference_path(case, stem).resolve()
               for case, stem, _ in CASES + PENDING}

    assert on_disk == claimed, f'orphaned: {sorted(on_disk - claimed)}'


def test_generating_the_same_document_twice_says_the_same_thing():
    """Self-test of the comparison, including the digit collapsing: if this
    fails the harness is measuring its own noise rather than the translation."""
    case, stem, template = CASES[0]
    work = prepare(case)

    generate(work, f'{stem}.yaml', template, 'once.docx')
    generate(work, f'{stem}.yaml', template, 'twice.docx')

    assert spoken(work / 'once.docx') == spoken(work / 'twice.docx')


def test_the_difference_report_notices_what_it_should():
    assert 'first difference at line 1' in difference(['a', 'b'], ['a', 'c'])
    assert 'more line(s)' in difference(['a', 'b'], ['a'])
    assert difference(['a'], ['a']) == 'no difference'


# ---------------------------------------------------------- the comparison

@pytest.mark.parametrize('case, stem, template', CASES)
def test_the_yaml_document_says_what_the_rdf_document_said(case, stem, template):
    """The document a .yaml produces says what its .rdf produced.

    This is what the whole migration was for, checked the only way that could
    have caught the clone-ordering defect: by reading the finished document.
    """
    expected, got, complaints = compare(case, stem, template)

    assert got == expected, (
        difference(expected, got)
        + chr(10) + 'the run said: ' + repr(complaints[:3]))


@pytest.mark.parametrize('case, stem, template', PENDING)
@pytest.mark.xfail(strict=True, reason=TEMPLATE_MISMATCH)
def test_the_case_whose_template_disagrees_with_the_ladder(case, stem, template):
    """Strict, so that renaming the tag in the template announces itself."""
    expected, got, complaints = compare(case, stem, template)

    assert got == expected, (
        difference(expected, got)
        + chr(10) + 'the run said: ' + repr(complaints[:3]))
